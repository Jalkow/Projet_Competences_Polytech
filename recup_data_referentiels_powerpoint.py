from pptx import Presentation
import re
import json

def cleanup_string(strig_to_cleanup):
    return re.sub(r"[^A-Za-z0-9!,./<>()'’àÀèÈéÉâÂêÊîÎôÔûÛœ\n -]+", "", strig_to_cleanup)

def delete_newline(string):
    return re.sub(r"\n", "", string)


TOLERANCE_ALIGNEMENT_HORIZONTAL = 200000 # En EMU (21 en pixels)

def extract_infos_competences(powerpoint_path):
    pres = Presentation(powerpoint_path)

    dict_infos = {}
    slide_number = 0
    slides = list(pres.slides)[2:]

    for slide in slides:
        slide_number += 1
        dict_infos[f"Comp{slide_number}"] = {"apprentissages_crit": {}}
        shape_situation_pro = None # La principale ou il y a marqué "Situations professionnelles :" (il peut en avoir une autre à droite car listing des situations coupé en deux)

        # Récup nom des différents niveau + setup pour récup les apprentissages critiques dans autre boucle
        for shape in slide.shapes:

            #Récupération nom niveau et setup
            for i in range(3, 6):
                if(f"niveau {i-2}" in shape.text.lower() or f"{i}A" in shape.text):
                    nom_app_crit = delete_newline(re.sub(r"[345]A|Niveau [123]", "", shape.text))
                    dict_infos[f"Comp{slide_number}"]["apprentissages_crit"][f"{i}A"] = {}
                    dict_infos[f"Comp{slide_number}"]["apprentissages_crit"][f"{i}A"]["nom"] = nom_app_crit


        for shape in slide.shapes:
            
            #Récupération nom compétence + numéro RNCP
            if("RNCP" in shape.text):
                dict_infos[f"Comp{slide_number}"]["RNCP"] = re.findall(r"RNCP[A-Z0-9]*", shape.text)
                dict_infos[f"Comp{slide_number}"]["nom_competence"] = delete_newline(cleanup_string(re.sub(r"RNCP[A-Z0-9]*", "", shape.text)))
            
            #Récupération composantes essentielles
            if("composantes essentielles" in shape.text.lower()):
                dict_infos[f"Comp{slide_number}"]["composantes_essentielles"] = cleanup_string(shape.text).splitlines()
                dict_infos[f"Comp{slide_number}"]["composantes_essentielles"].pop(0) # enlever le "Composantes essentielles :"
                dict_infos[f"Comp{slide_number}"]["composantes_essentielles"] = [composante for composante in dict_infos[f"Comp{slide_number}"]["composantes_essentielles"] if composante != ""] # Enlever les string vides de la liste


            # Récupération partie situations pro
            if("situations pro" in shape.text.lower()):
                shape_situation_pro = shape
                dict_infos[f"Comp{slide_number}"]["situations_pro"] = cleanup_string(shape.text).splitlines()
                dict_infos[f"Comp{slide_number}"]["situations_pro"].pop(0) # Enlever le "Situations professionnelles :"
                continue

            if(shape_situation_pro is not None):

                # Si aligné horizontalement avec la shape principale des situations pro
                if(abs(shape.top - shape_situation_pro.top) < TOLERANCE_ALIGNEMENT_HORIZONTAL):
                    for situation in cleanup_string(shape.text).splitlines():
                        dict_infos[f"Comp{slide_number}"]["situations_pro"].append(situation)

            # Récupération liste des apprentissages critiques par niveau
            if(shape_situation_pro is not None and shape.top >= shape_situation_pro.top + shape_situation_pro.height):
                # Si shape niveau de competence ou titre de partie -> skip
                if re.search(r"Niveau [123]|A[345]|APPRENTISSAGES CRITIQUES", shape.text):
                    continue
                
                for niveau in dict_infos[f"Comp{slide_number}"]["apprentissages_crit"].keys():
                    # Si on a déjà récup la liste d'un niveau, on passe au prochain niveau
                    if(dict_infos[f"Comp{slide_number}"]["apprentissages_crit"][niveau].get("app_list")):
                        continue

                    dict_infos[f"Comp{slide_number}"]["apprentissages_crit"][niveau]["app_list"] = cleanup_string(re.sub(r"AC[0-9]+.[0-9]+", "", shape.text)).splitlines()
                    break
    
    return dict_infos

    # sauvegarder sous json
    # with open(f"{nom_powerpoint_referentiel}_output.json", "w", encoding="utf-8") as f:
    #     json.dump(dict_infos, f, ensure_ascii=False, indent=4)


# extract_infos_competences("20251208_Referentiel_IDU.pptx")


